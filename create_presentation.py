"""
DebugMaster Presentation Generator - Premium Edition
Exciting, minimal slides with bold design
By Yasser Hamdan & Mohamad Dmayriye
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import os

# Colors - Neon Cyber Theme
DARK_BG = RGBColor(10, 10, 20)         # Deep dark
SURFACE = RGBColor(20, 20, 40)          # Card bg
PURPLE = RGBColor(138, 43, 226)         # Vibrant purple
NEON_GREEN = RGBColor(57, 255, 20)      # Neon green
NEON_PINK = RGBColor(255, 16, 240)      # Hot pink
CYAN = RGBColor(0, 255, 255)            # Cyan
GOLD = RGBColor(255, 215, 0)            # Gold
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(200, 200, 220)

def set_slide_background(slide, color):
    """Set slide background to solid color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_glow_rectangle(slide, left, top, width, height, color):
    """Add a glowing accent rectangle"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Make semi-transparent effect with lighter color
    return shape

def add_accent_line(slide, left, top, width, color):
    """Add a colored accent line"""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(6))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line

# ============================================
# SLIDE 1: EPIC TITLE
# ============================================
def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Decorative top accent
    add_accent_line(slide, Inches(0), Inches(0), Inches(13.33), PURPLE)

    # Bug emoji as logo
    logo = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(2), Inches(1.2))
    tf = logo.text_frame
    p = tf.paragraphs[0]
    p.text = "🐛"
    p.font.size = Pt(100)
    p.alignment = PP_ALIGN.CENTER

    # Main title with glow effect
    title = slide.shapes.add_textbox(Inches(0), Inches(2.5), Inches(13.33), Inches(1.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "DEBUG MASTER"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Tagline with neon effect
    tagline = slide.shapes.add_textbox(Inches(0), Inches(4), Inches(13.33), Inches(0.8))
    tf = tagline.text_frame
    p = tf.paragraphs[0]
    p.text = "LEARN TO DEBUG  •  LEVEL UP  •  COMPETE"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NEON_GREEN
    p.alignment = PP_ALIGN.CENTER

    # Decorative line
    add_accent_line(slide, Inches(4), Inches(5), Inches(5.33), CYAN)

    # Team names
    team = slide.shapes.add_textbox(Inches(0), Inches(5.8), Inches(13.33), Inches(0.8))
    tf = team.text_frame
    p = tf.paragraphs[0]
    p.text = "YASSER HAMDAN  &  MOHAMAD DMAYRIYE"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    # Bottom accent
    add_accent_line(slide, Inches(0), Inches(7.4), Inches(13.33), NEON_PINK)

# ============================================
# SLIDE 2: THE PROBLEM (Impactful)
# ============================================
def add_problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Big stat
    stat = slide.shapes.add_textbox(Inches(0), Inches(1), Inches(13.33), Inches(2))
    tf = stat.text_frame
    p = tf.paragraphs[0]
    p.text = "70%"
    p.font.size = Pt(150)
    p.font.bold = True
    p.font.color.rgb = NEON_PINK
    p.alignment = PP_ALIGN.CENTER

    # Stat explanation
    exp = slide.shapes.add_textbox(Inches(0), Inches(3.2), Inches(13.33), Inches(1))
    tf = exp.text_frame
    p = tf.paragraphs[0]
    p.text = "of developer time is spent DEBUGGING"
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Problem statement
    prob = slide.shapes.add_textbox(Inches(0), Inches(4.8), Inches(13.33), Inches(1.5))
    tf = prob.text_frame
    p = tf.paragraphs[0]
    p.text = "Yet learning to debug is..."
    p.font.size = Pt(28)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Pain points in boxes
    pains = [("😴", "BORING"), ("😤", "FRUSTRATING"), ("🚫", "NOT TAUGHT")]
    for i, (emoji, text) in enumerate(pains):
        x = Inches(1.5 + i * 3.8)

        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.5), Inches(3), Inches(1.3))
        box.fill.solid()
        box.fill.fore_color.rgb = SURFACE
        box.line.color.rgb = NEON_PINK
        box.line.width = Pt(3)

        # Emoji + text
        txt = slide.shapes.add_textbox(x, Inches(5.7), Inches(3), Inches(1))
        tf = txt.text_frame
        p = tf.paragraphs[0]
        p.text = f"{emoji} {text}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

# ============================================
# SLIDE 3: THE SOLUTION (Hero)
# ============================================
def add_solution_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # What if...
    intro = slide.shapes.add_textbox(Inches(0), Inches(0.8), Inches(13.33), Inches(0.8))
    tf = intro.text_frame
    p = tf.paragraphs[0]
    p.text = "What if debugging was..."
    p.font.size = Pt(32)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # FUN
    fun = slide.shapes.add_textbox(Inches(0), Inches(1.8), Inches(13.33), Inches(1.5))
    tf = fun.text_frame
    p = tf.paragraphs[0]
    p.text = "🎮 FUN"
    p.font.size = Pt(100)
    p.font.bold = True
    p.font.color.rgb = NEON_GREEN
    p.alignment = PP_ALIGN.CENTER

    # Feature cards
    features = [
        ("🐛", "FIX REAL BUGS", "Actual code execution", PURPLE),
        ("⚔️", "COMPETE", "Ranked Elo battles", NEON_PINK),
        ("🤖", "AI MENTOR", "Socratic guidance", CYAN),
        ("🏆", "LEVEL UP", "XP & achievements", GOLD)
    ]

    for i, (emoji, title, desc, color) in enumerate(features):
        x = Inches(0.8 + i * 3.1)

        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.8), Inches(2.8), Inches(2.8))
        card.fill.solid()
        card.fill.fore_color.rgb = SURFACE
        card.line.color.rgb = color
        card.line.width = Pt(4)

        # Emoji
        em = slide.shapes.add_textbox(x, Inches(4), Inches(2.8), Inches(1))
        tf = em.text_frame
        p = tf.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(50)
        p.alignment = PP_ALIGN.CENTER

        # Title
        t = slide.shapes.add_textbox(x, Inches(5), Inches(2.8), Inches(0.6))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        # Description
        d = slide.shapes.add_textbox(x, Inches(5.5), Inches(2.8), Inches(0.8))
        tf = d.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

# ============================================
# SLIDE 4: HOW IT WORKS (Gameplay)
# ============================================
def add_gameplay_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Title
    title = slide.shapes.add_textbox(Inches(0), Inches(0.3), Inches(13.33), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "⚡ HOW IT WORKS"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER

    # Steps with big numbers
    steps = [
        ("1", "SEE THE BUG", "Broken code + error message", "🔴"),
        ("2", "FIX IT", "Edit code in real-time", "✏️"),
        ("3", "RUN TESTS", "Janino compiler executes", "▶️"),
        ("4", "CELEBRATE!", "XP + confetti + level up", "🎉")
    ]

    for i, (num, title, desc, emoji) in enumerate(steps):
        y = Inches(1.5 + i * 1.4)

        # Big number
        num_box = slide.shapes.add_textbox(Inches(1), y, Inches(1.5), Inches(1.2))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = PURPLE
        p.alignment = PP_ALIGN.CENTER

        # Title
        t = slide.shapes.add_textbox(Inches(2.8), Inches(y.inches + 0.1), Inches(4), Inches(0.6))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Description
        d = slide.shapes.add_textbox(Inches(2.8), Inches(y.inches + 0.6), Inches(5), Inches(0.5))
        tf = d.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(18)
        p.font.color.rgb = LIGHT_GRAY

        # Emoji
        em = slide.shapes.add_textbox(Inches(10), y, Inches(2), Inches(1.2))
        tf = em.text_frame
        p = tf.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(60)
        p.alignment = PP_ALIGN.CENTER

    # Demo callout
    demo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(6.8), Inches(5.33), Inches(0.7))
    demo.fill.solid()
    demo.fill.fore_color.rgb = NEON_GREEN
    demo.line.fill.background()

    demo_txt = slide.shapes.add_textbox(Inches(4), Inches(6.9), Inches(5.33), Inches(0.5))
    tf = demo_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "🎬 LIVE DEMO"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BG
    p.alignment = PP_ALIGN.CENTER

# ============================================
# SLIDE 5: GAME MODES (Grid)
# ============================================
def add_modes_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Title
    title = slide.shapes.add_textbox(Inches(0), Inches(0.3), Inches(13.33), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "🎮 8 GAME MODES"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    # Mode grid (2 rows x 4 cols)
    modes = [
        ("⏱️", "Quick Fix", "60 sec per bug", NEON_GREEN),
        ("⚔️", "Battle Arena", "1v1 Elo ranked", NEON_PINK),
        ("🏃", "Speed Run", "Race the clock", CYAN),
        ("🧩", "Puzzle Mode", "No hints, 3x XP", PURPLE),
        ("📅", "Daily Challenge", "Fresh bugs daily", GOLD),
        ("❓", "Mystery Bug", "Random 5x XP", NEON_PINK),
        ("💀", "Survival", "1 life, endless", NEON_GREEN),
        ("📚", "Tutorial", "Learn step by step", CYAN)
    ]

    for i, (emoji, name, desc, color) in enumerate(modes):
        row = i // 4
        col = i % 4
        x = Inches(0.6 + col * 3.15)
        y = Inches(1.5 + row * 2.8)

        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.9), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = SURFACE
        card.line.color.rgb = color
        card.line.width = Pt(3)

        # Emoji
        em = slide.shapes.add_textbox(x, Inches(y.inches + 0.2), Inches(2.9), Inches(0.9))
        tf = em.text_frame
        p = tf.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(48)
        p.alignment = PP_ALIGN.CENTER

        # Name
        n = slide.shapes.add_textbox(x, Inches(y.inches + 1.1), Inches(2.9), Inches(0.5))
        tf = n.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Desc
        d = slide.shapes.add_textbox(x, Inches(y.inches + 1.6), Inches(2.9), Inches(0.6))
        tf = d.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

# ============================================
# SLIDE 6: RANKED SYSTEM (Epic)
# ============================================
def add_ranked_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Title
    title = slide.shapes.add_textbox(Inches(0), Inches(0.3), Inches(13.33), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "⚔️ RANKED BATTLES"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = NEON_PINK
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub = slide.shapes.add_textbox(Inches(0), Inches(1.1), Inches(13.33), Inches(0.6))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "Elo Rating System • Like Chess & League of Legends"
    p.font.size = Pt(22)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Tiers as badges
    tiers = [
        ("⚪", "Unranked", LIGHT_GRAY),
        ("🥉", "Bronze", RGBColor(205, 127, 50)),
        ("🥈", "Silver", RGBColor(192, 192, 192)),
        ("🥇", "Gold", GOLD),
        ("💎", "Diamond", CYAN),
        ("👑", "Master", PURPLE),
        ("🏆", "Legend", NEON_PINK)
    ]

    for i, (emoji, name, color) in enumerate(tiers):
        x = Inches(0.5 + i * 1.8)

        # Badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2), Inches(1.6), Inches(1.8))
        badge.fill.solid()
        badge.fill.fore_color.rgb = SURFACE
        badge.line.color.rgb = color
        badge.line.width = Pt(3)

        # Emoji
        em = slide.shapes.add_textbox(x, Inches(2.15), Inches(1.6), Inches(0.9))
        tf = em.text_frame
        p = tf.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(44)
        p.alignment = PP_ALIGN.CENTER

        # Name
        n = slide.shapes.add_textbox(x, Inches(3), Inches(1.6), Inches(0.6))
        tf = n.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

    # Key features
    features = [
        "🎯 Same bug, race to fix first",
        "📈 Win = gain Elo, Lose = drop",
        "🔥 Win streaks = bonus multiplier",
        "🏅 Seasonal rewards for top players"
    ]

    for i, feat in enumerate(features):
        row = i // 2
        col = i % 2
        x = Inches(1.5 + col * 5.5)
        y = Inches(4.5 + row * 0.9)

        f = slide.shapes.add_textbox(x, y, Inches(5), Inches(0.7))
        tf = f.text_frame
        p = tf.paragraphs[0]
        p.text = feat
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE

# ============================================
# SLIDE 7: TECH STACK (Minimal)
# ============================================
def add_tech_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Title
    title = slide.shapes.add_textbox(Inches(0), Inches(0.3), Inches(13.33), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "🛠️ TECH STACK"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER

    # Architecture visual
    arch_layers = [
        ("📱 UI", "Fragments + Material 3", PURPLE),
        ("🧠 Logic", "MVVM + ViewModels", NEON_GREEN),
        ("📦 Data", "Room DB (11 tables)", CYAN),
        ("☁️ Cloud", "Firebase + GPT-4", GOLD)
    ]

    for i, (layer, desc, color) in enumerate(arch_layers):
        y = Inches(1.5 + i * 1.3)

        # Layer bar
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), y, Inches(9.33), Inches(1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = SURFACE
        bar.line.color.rgb = color
        bar.line.width = Pt(4)

        # Layer name
        n = slide.shapes.add_textbox(Inches(2.3), Inches(y.inches + 0.2), Inches(3), Inches(0.6))
        tf = n.text_frame
        p = tf.paragraphs[0]
        p.text = layer
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = color

        # Description
        d = slide.shapes.add_textbox(Inches(6), Inches(y.inches + 0.25), Inches(5), Inches(0.5))
        tf = d.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE

    # Key tech badges
    techs = ["Java", "Hilt DI", "Janino", "Firebase", "GPT-4"]
    for i, tech in enumerate(techs):
        x = Inches(1.5 + i * 2.2)

        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(6.5), Inches(2), Inches(0.7))
        badge.fill.solid()
        badge.fill.fore_color.rgb = PURPLE
        badge.line.fill.background()

        t = slide.shapes.add_textbox(x, Inches(6.6), Inches(2), Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = tech
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

# ============================================
# SLIDE 8: THANK YOU (Finale)
# ============================================
def add_finale_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Top accent
    add_accent_line(slide, Inches(0), Inches(0), Inches(13.33), NEON_GREEN)

    # Big emoji
    em = slide.shapes.add_textbox(Inches(0), Inches(1), Inches(13.33), Inches(1.5))
    tf = em.text_frame
    p = tf.paragraphs[0]
    p.text = "🐛🎮🏆"
    p.font.size = Pt(80)
    p.alignment = PP_ALIGN.CENTER

    # Thank you
    thanks = slide.shapes.add_textbox(Inches(0), Inches(2.5), Inches(13.33), Inches(1.2))
    tf = thanks.text_frame
    p = tf.paragraphs[0]
    p.text = "THANK YOU!"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # App name
    app = slide.shapes.add_textbox(Inches(0), Inches(3.8), Inches(13.33), Inches(0.8))
    tf = app.text_frame
    p = tf.paragraphs[0]
    p.text = "DEBUG MASTER"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = NEON_GREEN
    p.alignment = PP_ALIGN.CENTER

    # Tagline
    tag = slide.shapes.add_textbox(Inches(0), Inches(4.6), Inches(13.33), Inches(0.6))
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.text = "Learn to Debug  •  Level Up  •  Compete"
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Questions box
    q_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(5.5), Inches(4.33), Inches(1))
    q_box.fill.solid()
    q_box.fill.fore_color.rgb = NEON_PINK
    q_box.line.fill.background()

    q = slide.shapes.add_textbox(Inches(4.5), Inches(5.7), Inches(4.33), Inches(0.6))
    tf = q.text_frame
    p = tf.paragraphs[0]
    p.text = "QUESTIONS?"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Team names
    team = slide.shapes.add_textbox(Inches(0), Inches(6.8), Inches(13.33), Inches(0.6))
    tf = team.text_frame
    p = tf.paragraphs[0]
    p.text = "Yasser Hamdan  &  Mohamad Dmayriye"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    # Bottom accent
    add_accent_line(slide, Inches(0), Inches(7.4), Inches(13.33), PURPLE)

def main():
    """Generate the premium presentation"""
    print("Creating DebugMaster Premium Presentation...")

    # Create presentation with 16:9 widescreen
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Add all slides
    add_title_slide(prs)       # 1. Title
    add_problem_slide(prs)     # 2. The Problem (70% stat)
    add_solution_slide(prs)    # 3. Solution (Feature cards)
    add_gameplay_slide(prs)    # 4. How it Works
    add_modes_slide(prs)       # 5. 8 Game Modes
    add_ranked_slide(prs)      # 6. Ranked Battles
    add_tech_slide(prs)        # 7. Tech Stack
    add_finale_slide(prs)      # 8. Thank You

    # Save
    output_path = os.path.join(os.path.dirname(__file__), "DebugMaster_Premium.pptx")
    prs.save(output_path)

    print("Presentation saved to:", output_path)
    print("8 slides with premium cyber design!")
    print("")
    print("TIP: Add animations in PowerPoint:")
    print("   - Select all objects on a slide")
    print("   - Animations > Fade or Fly In")
    print("   - Set to On Click or After Previous")

if __name__ == "__main__":
    main()
